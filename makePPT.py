from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import glob
import datetime as dt
from natsort import natsorted, ns
import sys
import pdb
prs = Presentation()

lyt=prs.slide_layouts[0] # choosing a slide layout
slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title

subtitle=slide.placeholders[1] # placeholder for subtitle

ID=sys.argv[1]
date=sys.argv[2]
CWID=sys.argv[3]

title.text=f'{CWID}' # title
subtitle.text=f'{date}' # subtitle

slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(100)

prs.slide_width = Inches(12.4)
prs.slide_height = Inches(11.1)

base_path = '/ix1/tibrahim/shared/tibrahim_jkofler/03-PMB/'
T1_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_T1/*.png'))
T1_anno_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_T1_anno/*.png'))

T2_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_T2/*.png'))
T2_anno_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_T2_anno/*.png'))

GRE_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_GRE/*.png'))
GRE_anno_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_GRE_anno/*.png'))

cam_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/resizedCam_adjusted/*.png'))
cam_anno_path = natsorted(glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/resizedCam_adjusted_anno/*.png'))

ap_layout = glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/ap/*.png')

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.add_picture(ap_layout[0], Inches(2.9), Inches(2.81))

for ind in range(len(T1_path)):
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(T1_path[ind], Inches(0.1), Inches(0.1))
    slide.shapes.add_picture(T2_path[ind], Inches(3.1), Inches(0.1))
    slide.shapes.add_picture(GRE_path[ind], Inches(6.1), Inches(0.1))
    slide.shapes.add_picture(cam_path[ind], Inches(9.1), Inches(0.1))
    
    slide.shapes.add_picture(T1_anno_path[ind], Inches(0.1), Inches(5.5))
    slide.shapes.add_picture(T2_anno_path[ind], Inches(3.1), Inches(5.5))
    slide.shapes.add_picture(GRE_anno_path[ind], Inches(6.1), Inches(5.5))
    slide.shapes.add_picture(cam_anno_path[ind], Inches(9.1), Inches(5.5))
# pdb.set_trace()
prs.save(f"{base_path}/{date}/{CWID}.pptx") # saving file
